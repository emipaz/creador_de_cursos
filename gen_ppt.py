import os
import json
import markdown
from pptx import Presentation
from mdx_math import MathExtension
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from prompt import get_full_html
from prompt import SYSTEM_MD, SYSTEM_PPT
from dotenv import load_dotenv , find_dotenv
if load_dotenv(find_dotenv()):
    import openai
    cliente = openai.Client()
else:
    print("Debe configurar una api_key de oopenai")


emb = OpenAIEmbeddings()

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size = 2000,
    chunk_overlap = 250
)

def base(curso, index, documentos):
    path = os.path.join(curso,f"base_{index}")
    bs = Chroma(embedding_function = emb, persist_directory = path)
    fragmentos = text_splitter.split_documents(documents=documentos)
    ids = [ d.metadata["id"] + "-frag-" + str(idx) for idx , d in enumerate(fragmentos,1)]
    return bs.add_documents(documents=fragmentos, ids=ids)

def markdown_to_html_with_math(markdown_text):
    # Configurar las extensiones
    md = markdown.Markdown(extensions=['extra', MathExtension(enable_dollar_delimiter=True)])
    
    # Convertir Markdown a HTML
    html = md.convert(markdown_text)
    return html

def guardar_html(html,curso, modulo, index):
    des_html = os.path.join(curso, modulo + ".html")
    full_html = get_full_html(html, modulo, index)
    with open(des_html, "w", encoding= "utf-8") as f:
        f.write(full_html)
    print("Se guardo en : ", des_html)

def crear_resumen(texto :str ) -> str:
    mensaje = [
        {
            "role": "system",
            "content": SYSTEM_MD},
        {
            "role": "user",
            "content": f"Trascripción del video : {texto}"
        }
    ]
    chat = cliente.chat.completions.create(
            model = "gpt-4o-mini", # modelo a usar
            messages = mensaje,    # lista de mensajes
            temperature = 0,)
    return chat.choices[0].message.content

def crear_slide(mark : str, ) -> str:
    mensaje = [
        {
            "role": "system",
            "content": SYSTEM_PPT},
        {
            "role": "user",
            "content": f"Convierte el siguiente markdown : {mark} eb una presentacion"
        }
    ]
    chat = cliente.chat.completions.create(
            model = "gpt-4o-mini", # modelo a usar
            response_format={ "type": "json_object" },
            messages = mensaje,    #t lista de mensajes
            temperature = 0,)

    return json.loads(chat.choices[0].message.content)

def crear_ppt(material, origen, modulo):
    curso = origen.split("\\")[-1]
    print(curso)
    titulo    = curso
    subtitulo = modulo
    slides    = material.get("slides", [] )
    nombre = curso + "_" +modulo +".ppt"
    path = os.path.join( origen ,nombre)
    prs = Presentation()
    
    # Diapositiva 1: Título
    slide_layout = prs.slide_layouts[0]  # Diseño de diapositiva de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titulo
    subtitle.text = subtitulo

    for slide in slides:
        titulo    = slide.get("titulo", "")
        # subtitulo = slide.get("subtitulo", "")
        contenido = slide.get("contenido", [] )
        hoja      = prs.slides.add_slide(slide_layout)
        title     = hoja.shapes.title
        content   = hoja.shapes.placeholders[1]
        title.text = titulo
        
        if isinstance (contenido, list):
            contenido = [str(x) for x in contenido]
            texto = "\n".join(contenido)
        else:
            texto = contenido
        content.text = texto

    prs.save(path)
    print("presentacion guardada en :", path)